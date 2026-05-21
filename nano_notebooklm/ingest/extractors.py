"""Document text extraction for PDF, PPTX, DOCX, and Markdown."""

from __future__ import annotations

from pathlib import Path

from nano_notebooklm.types import FileType, PageInfo


def extract_pdf(filepath: str) -> list[PageInfo]:
    """Extract text page-by-page from a PDF.

    fix-all v3 #L2: doc.close() previously ran only after the page loop
    completed; a malformed PDF that raised inside `page.get_text()`
    leaked the open document handle. Wrap in try/finally so the handle
    is always released.
    """
    import fitz

    doc = fitz.open(filepath)
    pages = []
    try:
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text().strip()
            if text and len(text) > 30:
                pages.append(PageInfo(
                    text=text,
                    page=page_num + 1,
                    total_pages=len(doc),
                ))
    finally:
        doc.close()
    return pages


def extract_pptx(filepath: str) -> list[PageInfo]:
    """Extract text slide-by-slide from a PPTX.

    fix-all v2 LOW F2: python-pptx uses lxml without disabling external
    entity processing, so a malicious .pptx (billion-laughs / XXE) can
    pin CPU or memory during the underlying XML parse. M4 already
    neutralised the page-count path (`_scan_file_pages`) by skipping
    python-pptx entirely there, but the real extraction path was left
    open. Defense: pre-check the .pptx zip envelope before letting
    python-pptx near the XML. Real decks have ~5-300 members totalling
    a few MB to ~100MB uncompressed; billion-laughs payloads inflate to
    GBs and/or carry thousands of tiny entity-defining members. Rejecting
    pathological envelopes upfront keeps the lxml call contained without
    regressing on real uploads. The zip-bomb thresholds are conservative
    enough that we'd expect zero false positives on classroom material.
    """
    import zipfile

    _PPTX_MAX_UNCOMPRESSED = 500 * 1024 * 1024  # 500MB
    _PPTX_MAX_MEMBERS = 5000
    _PPTX_MAX_INFLATION_RATIO = 200  # compressed → uncompressed

    try:
        with zipfile.ZipFile(filepath) as zf:
            infos = zf.infolist()
            if len(infos) > _PPTX_MAX_MEMBERS:
                raise ValueError(
                    f"pptx envelope rejected: {len(infos)} members exceeds "
                    f"safety cap {_PPTX_MAX_MEMBERS} (possible zip bomb)"
                )
            total_uncompressed = sum(i.file_size for i in infos)
            total_compressed = sum(i.compress_size for i in infos) or 1
            if total_uncompressed > _PPTX_MAX_UNCOMPRESSED:
                raise ValueError(
                    f"pptx envelope rejected: uncompressed {total_uncompressed} "
                    f"exceeds safety cap {_PPTX_MAX_UNCOMPRESSED}"
                )
            if total_uncompressed // total_compressed > _PPTX_MAX_INFLATION_RATIO:
                raise ValueError(
                    f"pptx envelope rejected: inflation ratio "
                    f"{total_uncompressed // total_compressed} exceeds "
                    f"safety cap {_PPTX_MAX_INFLATION_RATIO} (possible zip bomb)"
                )
    except zipfile.BadZipFile as exc:
        raise ValueError(f"pptx envelope rejected: not a valid zip ({exc})") from exc

    from pptx import Presentation

    prs = Presentation(filepath)
    slides = []
    for slide_num, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
        text = "\n".join(parts)
        if text and len(text) > 30:
            slides.append(PageInfo(
                text=text,
                slide=slide_num,
                total_slides=len(prs.slides),
            ))
    return slides


def extract_docx(filepath: str) -> list[PageInfo]:
    """Extract text paragraph-by-paragraph from a DOCX."""
    from docx import Document as DocxDocument

    doc = DocxDocument(filepath)
    sections = []
    current_heading = "Introduction"
    current_text: list[str] = []
    para_start = 1

    for i, para in enumerate(doc.paragraphs, 1):
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            # Save previous section
            text = "\n".join(current_text).strip()
            if text and len(text) > 30:
                sections.append(PageInfo(
                    text=text,
                    section=current_heading,
                    line_start=para_start,
                    line_end=i - 1,
                ))
            current_heading = para.text.strip() or current_heading
            current_text = []
            para_start = i
        else:
            if para.text.strip():
                current_text.append(para.text.strip())

    # Last section
    text = "\n".join(current_text).strip()
    if text and len(text) > 30:
        sections.append(PageInfo(
            text=text,
            section=current_heading,
            line_start=para_start,
            line_end=len(doc.paragraphs),
        ))
    return sections


def extract_markdown(filepath: str) -> list[PageInfo]:
    """Extract text section-by-section from a Markdown file."""
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()

    sections = []
    current_heading = "Introduction"
    current_text: list[str] = []
    line_start = 1

    for i, line in enumerate(content.split("\n"), 1):
        if line.startswith("#"):
            text = "\n".join(current_text).strip()
            if text and len(text) > 30:
                sections.append(PageInfo(
                    text=text,
                    section=current_heading,
                    line_start=line_start,
                    line_end=i - 1,
                ))
            current_heading = line.lstrip("#").strip()
            current_text = []
            line_start = i
        else:
            current_text.append(line)

    # Last section
    text = "\n".join(current_text).strip()
    if text and len(text) > 30:
        sections.append(PageInfo(
            text=text,
            section=current_heading,
            line_start=line_start,
            line_end=i if 'i' in dir() else line_start,
        ))
    return sections


def extract_file(
    filepath: str | Path,
    engine: str = "pymupdf",
    lang: str = "ch",
) -> tuple[list[PageInfo], FileType]:
    """Auto-detect file type and extract text.

    Args:
      engine: `pymupdf` (default, fast, drops formulae and tables) or
        `mineru` (slow ~10s/page, preserves LaTeX equations + HTML tables
        + extracted images). Only honoured for PDFs — other file types
        always fall back to their native extractor.
      lang: passed through to mineru when engine='mineru'. `ch` for
        Chinese, `en` for English.
    """
    filepath = Path(filepath)
    suffix = filepath.suffix.lower()

    if suffix == ".pdf" and engine == "mineru":
        # Lazy import to avoid loading torch / mineru deps in default path.
        from nano_notebooklm.ingest.extractors_mineru import extract_pdf_mineru

        return extract_pdf_mineru(filepath, lang=lang), FileType.PDF

    extractors = {
        ".pdf": (extract_pdf, FileType.PDF),
        ".pptx": (extract_pptx, FileType.PPTX),
        ".ppt": (extract_pptx, FileType.PPTX),
        ".docx": (extract_docx, FileType.DOCX),
        ".md": (extract_markdown, FileType.MARKDOWN),
        ".markdown": (extract_markdown, FileType.MARKDOWN),
        ".txt": (extract_markdown, FileType.TXT),
    }

    if suffix not in extractors:
        raise ValueError(f"Unsupported file type: {suffix}")

    func, file_type = extractors[suffix]
    return func(str(filepath)), file_type


# Directories to skip during recursive scanning
SKIP_DIRS = {".venv", "__pycache__", "node_modules", ".git", "docs"}
SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".ppt", ".docx", ".md", ".markdown", ".txt"}


def collect_files(directory: str | Path) -> list[Path]:
    """Recursively collect all supported files from a directory.

    R5-2 fix-all v4 dedup pass: when both `<stem>.pptx` (or `.ppt`) AND
    `<stem>.pdf` live in the same parent directory we keep ONLY the pptx
    and drop the pdf. Rationale:

    - PPTX carries semantically richer content (speaker notes, structure,
      slide titles) — better signal for chunking + KG extraction.
    - Users routinely export a pptx → pdf and drag both into the upload
      dialog. Pre-fix this indexed every slide twice (once via python-
      pptx, once via pdfplumber), inflating chunks.json, the FAISS
      index, and the sources panel (two rows for "the same lecture").
    - The Reader still gets a PDF view: pptx_pdf.convert_pptx_to_pdf
      renders a sidecar at upload time, served via
      `_resolve_pptx_pdf_sidecar` — the user-uploaded pdf is redundant.

    The drop is silent at this layer (chunker-level) but the upload
    pipeline can surface counts via `len(collect_files(upload_dir))` vs
    file count if it wants to warn the user.
    """
    directory = Path(directory)
    files = []
    for ext in SUPPORTED_EXTENSIONS:
        for f in directory.rglob(f"*{ext}"):
            if not any(skip in f.parts for skip in SKIP_DIRS):
                files.append(f)
    return _dedupe_pptx_pdf_pairs(sorted(files))


def _dedupe_pptx_pdf_pairs(files: list[Path]) -> list[Path]:
    """Drop a `<stem>.pdf` when a sibling `<stem>.pptx` (or `.ppt`) exists.

    Sibling = same parent directory + same stem (case-sensitive on the
    file system the caller hands us). Comparison is exact-stem only; if
    a user wants both `lecture1.pptx` and `lecture1-handout.pdf` to be
    indexed, they remain (different stems).
    """
    by_dir_stem_pptx: set[tuple[Path, str]] = set()
    for f in files:
        if f.suffix.lower() in (".pptx", ".ppt"):
            by_dir_stem_pptx.add((f.parent, f.stem))
    if not by_dir_stem_pptx:
        return files
    kept = []
    for f in files:
        if f.suffix.lower() == ".pdf" and (f.parent, f.stem) in by_dir_stem_pptx:
            continue
        kept.append(f)
    return kept
